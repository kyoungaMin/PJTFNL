"""
고객사 제안서 PPT 생성 스크립트
— AI 기반 수요 변동성 최적화 SaaS 제안서
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# ─── 디자인 상수 ─────────────────────────────────
NAVY      = RGBColor(0x0D, 0x1B, 0x2A)   # 진한 네이비
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
MID_BLUE  = RGBColor(0x2E, 0x86, 0xAB)
ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # 밝은 청록
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xF0, 0xF0, 0xF5)
GRAY      = RGBColor(0x6C, 0x75, 0x7D)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
RED       = RGBColor(0xE0, 0x4F, 0x5F)
GREEN     = RGBColor(0x2D, 0xCE, 0x89)
ORANGE    = RGBColor(0xFB, 0x6D, 0x48)
YELLOW    = RGBColor(0xFF, 0xC1, 0x07)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color):
    """슬라이드 배경색 설정"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    """사각형 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        elem = shape.fill._fill
        import lxml.etree as etree
        srgb = elem.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            a = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            a.set('val', str(int(alpha * 1000)))
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="맑은 고딕"):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=14, color=DARK_GRAY, bold=False,
             alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4),
             font_name="맑은 고딕", level=0):
    """텍스트 프레임에 문단 추가"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


def add_icon_card(slide, left, top, width, height, icon_text, title, desc,
                  icon_color=ACCENT, bg_color=WHITE):
    """아이콘 + 제목 + 설명 카드"""
    # 배경
    card = add_rect(slide, left, top, width, height, bg_color)
    # 상단 컬러 바
    add_rect(slide, left, top, width, Inches(0.06), icon_color)
    # 아이콘 (텍스트 기반)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2),
                 Inches(0.6), Inches(0.6), icon_text,
                 font_size=28, color=icon_color, bold=True, alignment=PP_ALIGN.CENTER)
    # 제목
    add_text_box(slide, left + Inches(0.3), top + Inches(0.75),
                 width - Inches(0.6), Inches(0.4), title,
                 font_size=16, color=DARK_GRAY, bold=True)
    # 설명
    add_text_box(slide, left + Inches(0.3), top + Inches(1.1),
                 width - Inches(0.6), height - Inches(1.3), desc,
                 font_size=11, color=GRAY)


def add_table_slide(slide, left, top, rows_data, col_widths, header_color=NAVY):
    """테이블 추가"""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 0
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                          sum(col_widths), Inches(0.4) * n_rows)
    table = table_shape.table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.name = "맑은 고딕"

            if ri == 0:  # header
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                para.font.color.rgb = WHITE
                para.font.bold = True
                para.alignment = PP_ALIGN.CENTER
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else LIGHT_GRAY
                para.font.color.rgb = DARK_GRAY
                para.alignment = PP_ALIGN.LEFT

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def build_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ============================================================
    # SLIDE 1: 표지
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, NAVY)
    # 상단 장식 바
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
    # 좌측 액센트 라인
    add_rect(slide, Inches(0.8), Inches(2.0), Inches(0.08), Inches(3.5), ACCENT)
    # 메인 타이틀
    add_text_box(slide, Inches(1.2), Inches(2.0), Inches(8), Inches(0.8),
                 "AI 기반 수요 변동성 최적화 SaaS",
                 font_size=38, color=WHITE, bold=True)
    # 서브 타이틀
    add_text_box(slide, Inches(1.2), Inches(2.8), Inches(10), Inches(0.7),
                 "반도체 부품·소재 제조기업을 위한 수요예측 및 재고 리스크 최적화 플랫폼",
                 font_size=20, color=ACCENT)
    # 구분선
    add_rect(slide, Inches(1.2), Inches(3.7), Inches(4), Inches(0.03), ACCENT)
    # 부가 정보
    add_text_box(slide, Inches(1.2), Inches(4.0), Inches(8), Inches(0.5),
                 "ERP 위에 얹는 AI 의사결정 레이어 — 기존 데이터만으로 3개월 내 PoC 가능",
                 font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC))
    # 하단 정보
    add_text_box(slide, Inches(1.2), Inches(5.8), Inches(6), Inches(0.4),
                 "2026년 3월",
                 font_size=14, color=GRAY)
    # 하단 바
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), ACCENT)

    # ============================================================
    # SLIDE 2: 목차
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), Inches(0.5), SLIDE_H, NAVY)
    add_text_box(slide, Inches(1.0), Inches(0.5), Inches(5), Inches(0.7),
                 "Contents", font_size=32, color=NAVY, bold=True)
    add_rect(slide, Inches(1.0), Inches(1.2), Inches(2), Inches(0.04), ACCENT)

    toc_items = [
        ("01", "현황 분석 — 생산기획의 한계"),
        ("02", "솔루션 개요 — AI 의사결정 레이어"),
        ("03", "핵심 기능 — 5가지 핵심 모듈"),
        ("04", "기술 아키텍처 — 시스템 구성"),
        ("05", "데이터 파이프라인 — 9단계 분석 엔진"),
        ("06", "방법론 — AI/ML 모델링 전략"),
        ("07", "검증 결과 — 실데이터 기반 성과"),
        ("08", "기대 효과 — 정량·정성 목표"),
        ("09", "구축 로드맵 — 추진 일정"),
        ("10", "Q&A"),
    ]
    for i, (num, title) in enumerate(toc_items):
        y = Inches(1.6) + Inches(0.52) * i
        add_text_box(slide, Inches(1.2), y, Inches(0.7), Inches(0.45),
                     num, font_size=20, color=ACCENT, bold=True)
        add_text_box(slide, Inches(2.0), y, Inches(8), Inches(0.45),
                     title, font_size=16, color=DARK_GRAY)

    # ============================================================
    # SLIDE 3: 현황 분석 — 생산기획의 한계
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.6),
                 "01  현황 분석 — 생산기획의 한계",
                 font_size=26, color=WHITE, bold=True)

    # 좌측: 현재 문제점
    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
                 "현재 ERP 기반 생산기획의 한계",
                 font_size=18, color=NAVY, bold=True)

    problems = [
        ("수주 급증  →  납기 지연", "수요 변동 사전 감지 불가, 항상 사후 대응"),
        ("수주 급감  →  재고 과잉", "감산 판단이 늦어 과잉 생산·재고 비용 증가"),
        ("엑셀 수작업 의존", "담당자 경험·감에 의존한 생산계획 수립"),
        ("외부 변수 미반영", "환율·반도체 시장·거시경제 변동 예측 반영 불가"),
    ]
    for i, (title, desc) in enumerate(problems):
        y = Inches(2.0) + Inches(0.7) * i
        add_rect(slide, Inches(0.8), y, Inches(0.08), Inches(0.5), RED)
        add_text_box(slide, Inches(1.1), y, Inches(5), Inches(0.3),
                     title, font_size=14, color=DARK_GRAY, bold=True)
        add_text_box(slide, Inches(1.1), y + Inches(0.28), Inches(5), Inches(0.3),
                     desc, font_size=11, color=GRAY)

    # 우측: 핵심 인사이트
    add_rect(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(4.8),
             LIGHT_GRAY)
    add_text_box(slide, Inches(7.4), Inches(1.6), Inches(4.8), Inches(0.5),
                 "핵심 인사이트",
                 font_size=18, color=NAVY, bold=True)
    txBox = add_text_box(slide, Inches(7.4), Inches(2.2), Inches(4.8), Inches(3.5),
                         "", font_size=14, color=DARK_GRAY)
    tf = txBox.text_frame
    tf.paragraphs[0].text = '"문제는 데이터가 없는 게 아니라,'
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "맑은 고딕"
    add_para(tf, '데이터를 미래 의사결정으로', font_size=16, color=NAVY, bold=True)
    add_para(tf, '연결하는 구조가 없다는 것"', font_size=16, color=NAVY, bold=True)
    add_para(tf, '', font_size=8)
    add_para(tf, '• ERP는 과거 기록 시스템', font_size=13, color=DARK_GRAY)
    add_para(tf, '• 수주 조회는 가능하지만 예측은 불가', font_size=13, color=DARK_GRAY)
    add_para(tf, '• 재고 현황은 보이지만 리스크 감지는 불가', font_size=13, color=DARK_GRAY)
    add_para(tf, '• 생산 실적은 있지만 최적 생산량 산출은 불가', font_size=13, color=DARK_GRAY)

    # 하단 배너
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.8), SLIDE_W, Inches(0.8), NAVY)
    add_text_box(slide, Inches(0.8), SLIDE_H - Inches(0.65), Inches(11), Inches(0.5),
                 "본 솔루션은 ERP를 대체하는 것이 아닌, ERP 위에 얹는 AI 의사결정 레이어입니다",
                 font_size=15, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 4: 솔루션 개요
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
                 "02  솔루션 개요 — AI 의사결정 레이어",
                 font_size=26, color=WHITE, bold=True)

    # 3개 핵심 포지셔닝
    cards = [
        ("01", "ERP 보완", "ERP를 대체하지 않음\nERP 위에 얹는\nAI 의사결정 레이어",
         MID_BLUE),
        ("02", "특화 도구", "복잡한 전사 SCM이 아닌\n생산기획 의사결정\n특화 AI 도구",
         ACCENT),
        ("03", "빠른 도입", "기존 ERP 데이터만으로\n시작 가능\n3개월 내 PoC 완료",
         GREEN),
    ]
    for i, (num, title, desc, color) in enumerate(cards):
        x = Inches(0.8) + Inches(4.1) * i
        add_rect(slide, x, Inches(1.5), Inches(3.7), Inches(2.5), color)
        add_text_box(slide, x + Inches(0.3), Inches(1.7), Inches(3.1), Inches(0.4),
                     num, font_size=32, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.3), Inches(2.15), Inches(3.1), Inches(0.4),
                     title, font_size=20, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.3), Inches(2.7), Inches(3.1), Inches(1.2),
                     desc, font_size=13, color=WHITE)

    # 대상 조직
    add_text_box(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(0.4),
                 "대상 사용자", font_size=18, color=NAVY, bold=True)
    targets = [
        ("생산기획 팀", "주간/월간 생산계획 수립"),
        ("운영/SCM 담당", "자재 조달·재고 관리"),
        ("영업 관리", "고객사별 수주 동향 파악"),
        ("경영진", "데이터 기반 의사결정 리포트"),
    ]
    for i, (role, desc) in enumerate(targets):
        x = Inches(0.8) + Inches(3.0) * i
        add_rect(slide, x, Inches(4.8), Inches(2.7), Inches(0.35), LIGHT_GRAY)
        add_text_box(slide, x + Inches(0.15), Inches(4.82), Inches(1.2), Inches(0.3),
                     role, font_size=12, color=NAVY, bold=True)
        add_text_box(slide, x + Inches(1.4), Inches(4.82), Inches(1.2), Inches(0.3),
                     desc, font_size=10, color=GRAY)

    # 적합 기업
    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(5), Inches(0.4),
                 "이런 기업에 적합합니다", font_size=18, color=NAVY, bold=True)
    fits = [
        "고객사별 주문 변동 폭이 큰 기업",
        "매주 생산계획을 수정하는 기업",
        "매출-생산 간 괴리가 자주 발생하는 기업",
        "반도체·부품·소재 등 글로벌 지표 영향이 큰 산업",
    ]
    txBox = add_text_box(slide, Inches(0.8), Inches(5.9), Inches(11), Inches(1.2),
                         "", font_size=13, color=DARK_GRAY)
    for fit in fits:
        add_para(txBox.text_frame, f"   {fit}", font_size=13, color=DARK_GRAY,
                 space_before=Pt(2), space_after=Pt(2))

    # ============================================================
    # SLIDE 5: 핵심 기능 — 5가지 모듈
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
                 "03  핵심 기능 — 5가지 핵심 모듈",
                 font_size=26, color=WHITE, bold=True)

    modules = [
        ("01", "AI 수요 예측",
         "제품별 확률 밴드 예측\n• P10(낙관) / P50(중앙) / P90(비관)\n• 주간: 1w/2w/4w\n• 월간: 1m/3m/6m\n• LightGBM Quantile Regression",
         ACCENT),
        ("02", "외부 변수 반영",
         "30+개 경제·산업 지표 반영\n• 환율(4통화), SOX지수, DRAM현물가\n• FRED/EIA/ECOS/관세청 API\n• 거시경제·무역·원자재 자동 수집",
         MID_BLUE),
        ("03", "리스크 스코어링",
         "4유형 자동 평가 + A~F 등급\n• 결품 리스크 (35%)\n• 과잉 리스크 (25%)\n• 납기 리스크 (25%)\n• 마진 리스크 (15%)",
         DARK_BLUE),
        ("04", "생산 최적화",
         "예측+재고+캐파 기반 최적 생산량\n• 리스크별 동적 조정\n• 결품 위험 시 P90 상향\n• 과잉 위험 시 감량\n• 긴급수주 자동 반영",
         RGBColor(0x6C, 0x5C, 0xE7)),
        ("05", "발주 최적화",
         "BOM 전개 + EOQ/ROP 기반 추천\n• 공급사 가중점수 추천\n• 리드타임·단가·신뢰도 종합\n• 대체 공급사 자동 제시\n• 발주 일정·긴급도 산출",
         RGBColor(0xE1, 0x7C, 0x05)),
    ]
    for i, (num, title, desc, color) in enumerate(modules):
        x = Inches(0.4) + Inches(2.5) * i
        # 카드 배경
        add_rect(slide, x, Inches(1.4), Inches(2.3), Inches(5.3), LIGHT_GRAY)
        # 상단 컬러 헤더
        add_rect(slide, x, Inches(1.4), Inches(2.3), Inches(1.0), color)
        add_text_box(slide, x + Inches(0.15), Inches(1.45), Inches(2.0), Inches(0.35),
                     num, font_size=22, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.15), Inches(1.85), Inches(2.0), Inches(0.4),
                     title, font_size=15, color=WHITE, bold=True)
        # 설명
        add_text_box(slide, x + Inches(0.15), Inches(2.6), Inches(2.0), Inches(3.8),
                     desc, font_size=10, color=DARK_GRAY)

    # ============================================================
    # SLIDE 6: 사용자 시나리오
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
                 "03+  사용자 시나리오 — 현업 활용 흐름",
                 font_size=26, color=WHITE, bold=True)

    steps = [
        ("STEP 1", "로그인\n대시보드 확인", MID_BLUE),
        ("STEP 2", "주별/월별\n예측 단위 선택", ACCENT),
        ("STEP 3", "수요 예측 확인\nP10/P50/P90 밴드", MID_BLUE),
        ("STEP 4", "변동성 큰 품목\n자동 식별 (C등급+)", ACCENT),
        ("STEP 5", "권장 조치 확인\n생산계획 수정", MID_BLUE),
        ("STEP 6", "예측 vs 실적\n모델 신뢰도 확인", ACCENT),
    ]
    for i, (step, desc, color) in enumerate(steps):
        x = Inches(0.5) + Inches(2.1) * i
        # 동그란 스텝 번호
        add_rect(slide, x, Inches(1.6), Inches(1.8), Inches(1.4), color)
        add_text_box(slide, x + Inches(0.15), Inches(1.65), Inches(1.5), Inches(0.35),
                     step, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), Inches(2.05), Inches(1.5), Inches(0.8),
                     desc, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)
        # 화살표
        if i < 5:
            add_text_box(slide, x + Inches(1.8), Inches(2.0), Inches(0.3), Inches(0.5),
                         "→", font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)

    # 하단: 핵심 가치
    add_rect(slide, Inches(0.5), Inches(3.5), Inches(12), Inches(0.06), ACCENT)
    add_text_box(slide, Inches(0.5), Inches(3.8), Inches(12), Inches(0.5),
                 '생산회의에서 "감"이 아닌 데이터 기반 설명이 가능해집니다',
                 font_size=18, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    # 조치 제안 테이블
    add_text_box(slide, Inches(0.8), Inches(4.6), Inches(5), Inches(0.4),
                 "상황별 자동 조치 제안", font_size=16, color=NAVY, bold=True)
    action_data = [
        ["상황", "자동 제안", "우선순위"],
        ["수요 급증 (결품 위험)", "긴급 발주 / 증산 검토", "Critical"],
        ["수요 급감 (과잉 위험)", "감산 / 발주 보류", "High"],
        ["납기 지연 위험", "대체 공급사 / 긴급 입고", "Critical"],
        ["변동성 과다", "안전재고 상향", "Medium"],
        ["고객 집중 리스크", "우선 대응 품목 알림", "High"],
    ]
    add_table_slide(slide, Inches(0.8), Inches(5.0), action_data,
                    [Inches(3.5), Inches(4.5), Inches(2.0)])

    # ============================================================
    # SLIDE 7: 기술 아키텍처
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
                 "04  기술 아키텍처 — 시스템 구성",
                 font_size=26, color=WHITE, bold=True)

    # 3-Tier 아키텍처 시각화
    tiers = [
        ("사용자 (브라우저)", "Next.js 대시보드  |  Vercel 배포  |  반응형 UI", ACCENT, Inches(1.4)),
        ("FastAPI 백엔드 서버", "예측 API  |  리스크 API  |  대시보드 API  |  Swagger 자동 문서", MID_BLUE, Inches(2.6)),
        ("PostgreSQL (Supabase)", "29개 테이블  |  REST API  |  실시간 구독  |  RLS 보안  |  RBAC 4단계", DARK_BLUE, Inches(3.8)),
        ("데이터 파이프라인 (Python)", "CSV 적재 → 외부 API 수집 → 9단계 분석 → 생산·발주 최적화", NAVY, Inches(5.0)),
    ]
    for title, desc, color, y in tiers:
        add_rect(slide, Inches(1.0), y, Inches(11), Inches(0.9), color)
        add_text_box(slide, Inches(1.3), y + Inches(0.05), Inches(4), Inches(0.4),
                     title, font_size=16, color=WHITE, bold=True)
        add_text_box(slide, Inches(1.3), y + Inches(0.45), Inches(10.5), Inches(0.4),
                     desc, font_size=12, color=WHITE)

    # 화살표
    for y in [Inches(2.3), Inches(3.5), Inches(4.7)]:
        add_text_box(slide, Inches(6.0), y, Inches(1), Inches(0.35),
                     "▼", font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

    # 외부 데이터 소스
    add_text_box(slide, Inches(1.0), Inches(6.1), Inches(3), Inches(0.35),
                 "외부 데이터 소스", font_size=14, color=NAVY, bold=True)
    sources = ["FRED (미연준)", "EIA (에너지청)", "관세청 API", "ECOS (한국은행)", "ERP CSV"]
    for i, src in enumerate(sources):
        x = Inches(1.0) + Inches(2.2) * i
        add_rect(slide, x, Inches(6.5), Inches(2.0), Inches(0.4), LIGHT_GRAY)
        add_text_box(slide, x + Inches(0.1), Inches(6.52), Inches(1.8), Inches(0.35),
                     src, font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 8: 데이터 파이프라인 — 9단계
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
                 "05  데이터 파이프라인 — 9단계 분석 엔진",
                 font_size=26, color=WHITE, bold=True)

    pipeline_data = [
        ["Step", "모듈명", "입력", "출력", "설명"],
        ["S0", "주별·월별 집계", "수주, 매출, 생산", "집계 4테이블", "ISO 주차 캘린더 + 다차원 집계"],
        ["S1", "일간 추정 재고", "재고, 생산, 매출", "daily_inventory", "월초 스냅샷 기반 일간 보간"],
        ["S2", "리드타임 통계", "구매발주", "product_lead_time", "AVG / P90 리드타임"],
        ["S3", "피처 엔지니어링", "ERP + 외부지표", "feature_store", "주간 46피처 / 월간 35피처"],
        ["S4", "수요예측 모델", "feature_store", "forecast_result", "LightGBM P10/P50/P90"],
        ["S5", "리스크 스코어링", "예측 + 재고", "risk_score", "4유형, A~F 등급"],
        ["S6", "조치 큐 생성", "risk_score", "action_queue", "자동 조치 제안"],
        ["S7", "생산 최적화", "예측+캐파+리스크", "production_plan", "최적 생산량 산출"],
        ["S8", "발주 최적화", "BOM+리드타임", "purchase_rec.", "BOM전개+EOQ+공급사추천"],
    ]
    add_table_slide(slide, Inches(0.3), Inches(1.3), pipeline_data,
                    [Inches(0.7), Inches(2.0), Inches(2.2), Inches(2.2), Inches(3.5)])

    # 하단 보충
    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.4),
                 "▶ 주간 파이프라인(S0~S8) + 월간 파이프라인(S3m~S4m)  |  자동화 스케줄링 지원  |  단일 CLI로 통합 실행",
                 font_size=13, color=MID_BLUE, bold=True)

    # DB 구조 요약
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(3), Inches(0.4),
                 "DB 구조: 29개 테이블", font_size=14, color=NAVY, bold=True)
    db_data = [
        ["구분", "테이블 수", "비고"],
        ["마스터", "3", "제품, 공급사, 고객사"],
        ["트랜잭션", "6", "수주, 매출, 생산, 발주, 재고, BOM"],
        ["외부지표", "3", "경제, 무역, 환율"],
        ["분석/ML", "8", "피처스토어, 예측, 리스크, 조치큐"],
        ["최적화", "2", "생산계획, 발주추천"],
        ["집계/인증", "7", "주별/월별 집계, 사용자"],
    ]
    add_table_slide(slide, Inches(6.5), Inches(6.2), db_data,
                    [Inches(1.3), Inches(1.0), Inches(3.5)])

    # ============================================================
    # SLIDE 9: 방법론 — AI/ML 모델링 전략
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
                 "06  방법론 — AI/ML 모델링 전략",
                 font_size=26, color=WHITE, bold=True)

    # 좌: 예측 모델
    add_rect(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5), LIGHT_GRAY)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
                 "수요예측 모델 — LightGBM Quantile Regression",
                 font_size=16, color=NAVY, bold=True)

    model_items = [
        ("모델 선정 사유", "• P10/P50/P90 분위 예측 직접 산출\n• 외부 피처(환율, 지수) 활용 용이\n• 범주형 변수 네이티브 처리\n• 빠른 학습 속도, 과적합 방지"),
        ("피처 엔지니어링 (10개 카테고리)", "A. 수주 이력 래그 (lag1~52주, 이동평균)\nB. 모멘텀 (변화율, 차분)\nC. 변동성 (표준편차, CV, 비영수주율)\nD. 공급측 (생산, 매출, 재고, BTB 비율)\nE. 고객 집중도 (HHI, TOP1/3 비율)\nF. 가격 (평균단가, 주문평균값)\nG. 반도체 시장 (SOX, DRAM, NAND)\nH. 환율 (USD/JPY/EUR/CNY→KRW)\nI. 거시경제 (금리, 산업생산, CPI, PMI)\nJ. 무역 (반도체 수출입, 무역수지)\nK. 시간 (주차, 월, 분기, 연말 여부)"),
        ("Fallback 전략", "• LightGBM 미설치 시 이동평균 자동 전환\n• 학습 데이터 부족 제품: 가중 이동평균 적용"),
    ]
    y_pos = Inches(2.0)
    for title, desc in model_items:
        add_text_box(slide, Inches(0.8), y_pos, Inches(5.2), Inches(0.3),
                     title, font_size=12, color=MID_BLUE, bold=True)
        tb = add_text_box(slide, Inches(0.8), y_pos + Inches(0.3), Inches(5.2), Inches(1.5),
                          desc, font_size=10, color=DARK_GRAY)
        lines = desc.count('\n') + 1
        y_pos += Inches(0.3) + Inches(0.18) * lines + Inches(0.1)

    # 우: 리스크 & 최적화
    add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5), LIGHT_GRAY)
    add_text_box(slide, Inches(7.1), Inches(1.5), Inches(5), Inches(0.4),
                 "리스크 스코어링 + 생산·발주 최적화",
                 font_size=16, color=NAVY, bold=True)

    right_items = [
        ("리스크 스코어링 (S5)", "• 결품(35%): P90수요 대비 재고 부족도\n• 과잉(25%): 재고일수 90일 초과 기준\n• 납기(25%): P90 리드타임 대비 잔여일\n• 마진(15%): 원가 상승·환율 변동 압박\n• 등급: A(0~20) ~ F(81~100)"),
        ("생산 최적화 (S7)", "• 수요예측 + 현재재고 + 안전재고 → 순소요량\n• 생산 캐파시티 제약 반영 (×1.2 버퍼)\n• 리스크 등급별 동적 조정\n  - D/F + 결품↑ → P90 기반 증산\n  - D/F + 과잉↑ → 10% 감량\n• 미처리 긴급수주 자동 반영"),
        ("발주 최적화 (S8)", "• BOM 전개 → 자재별 총소요량 산출\n• 안전재고·재주문점(ROP) 산출\n• EOQ + Lot-for-Lot 하이브리드\n• 공급사 가중점수 추천\n  (리드타임 40% + 단가 35% + 신뢰도 25%)\n• 대체 공급사 자동 제시"),
    ]
    y_pos = Inches(2.0)
    for title, desc in right_items:
        add_text_box(slide, Inches(7.1), y_pos, Inches(5.2), Inches(0.3),
                     title, font_size=12, color=MID_BLUE, bold=True)
        add_text_box(slide, Inches(7.1), y_pos + Inches(0.3), Inches(5.2), Inches(1.5),
                     desc, font_size=10, color=DARK_GRAY)
        lines = desc.count('\n') + 1
        y_pos += Inches(0.3) + Inches(0.18) * lines + Inches(0.1)

    # ============================================================
    # SLIDE 10: 검증 결과 — 실데이터 성과
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
                 "07  검증 결과 — 실데이터 기반 성과",
                 font_size=26, color=WHITE, bold=True)

    # 핵심 수치 카드 4개
    kpi_cards = [
        ("13,321", "리스크 분석 대상\n제품 수", MID_BLUE),
        ("3,562", "생산 계획\n산출 건수", ACCENT),
        ("1,356", "발주 추천\n산출 건수", GREEN),
        ("100%", "C등급 이상\n커버리지", RGBColor(0x6C, 0x5C, 0xE7)),
    ]
    for i, (value, label, color) in enumerate(kpi_cards):
        x = Inches(0.5) + Inches(3.15) * i
        add_rect(slide, x, Inches(1.3), Inches(2.9), Inches(1.4), color)
        add_text_box(slide, x + Inches(0.2), Inches(1.4), Inches(2.5), Inches(0.7),
                     value, font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(2.1), Inches(2.5), Inches(0.5),
                     label, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

    # 좌: 리스크 등급 분포
    add_text_box(slide, Inches(0.5), Inches(3.0), Inches(5), Inches(0.4),
                 "리스크 등급 분포 (검증 완료)", font_size=16, color=NAVY, bold=True)
    risk_data = [
        ["등급", "건수", "비율", "의미"],
        ["A (안전)", "10,199", "76.6%", "리스크 낮음"],
        ["B (주의)", "1,564", "11.7%", "모니터링 필요"],
        ["C (경고)", "1,512", "11.4%", "조치 검토 대상"],
        ["D (위험)", "46", "0.3%", "즉시 대응 필요"],
        ["F (심각)", "0", "0%", "-"],
    ]
    add_table_slide(slide, Inches(0.5), Inches(3.5), risk_data,
                    [Inches(1.2), Inches(1.0), Inches(0.8), Inches(2.5)])

    # 우: S7 생산 최적화 결과
    add_text_box(slide, Inches(6.5), Inches(3.0), Inches(5), Inches(0.4),
                 "S7 생산 최적화 결과", font_size=16, color=NAVY, bold=True)
    s7_data = [
        ["항목", "값"],
        ["총 산출 건수", "3,562건"],
        ["Critical 우선순위", "2,125건 (60%)"],
        ["계획 유형 - 증산", "981건"],
        ["계획 유형 - 감산", "824건"],
        ["계획 유형 - 유지", "452건"],
        ["계획 유형 - 신규", "1,305건"],
    ]
    add_table_slide(slide, Inches(6.5), Inches(3.5), s7_data,
                    [Inches(2.5), Inches(3.0)])

    # 하단: S8 발주 최적화 결과
    add_text_box(slide, Inches(0.5), Inches(5.8), Inches(5), Inches(0.4),
                 "S8 발주 최적화 결과", font_size=16, color=NAVY, bold=True)
    s8_data = [
        ["항목", "값", "비고"],
        ["총 추천 건수", "1,356건", "BOM 전개 기반"],
        ["긴급도 Critical", "1,220건 (90%)", "즉시 발주 필요"],
        ["발주방식 Lot-for-Lot", "1,218건 (90%)", "순소요량 기반"],
        ["발주방식 EOQ", "138건 (10%)", "경제적 주문량"],
        ["공급사 매핑율", "11.7%", "이력 기반, 확장 가능"],
    ]
    add_table_slide(slide, Inches(0.5), Inches(6.2), s8_data,
                    [Inches(2.2), Inches(1.8), Inches(2.5)])

    # 우하단: 무결성 검증
    add_text_box(slide, Inches(6.5), Inches(5.8), Inches(5), Inches(0.4),
                 "무결성 검증 결과", font_size=16, color=NAVY, bold=True)
    integrity_data = [
        ["검증 항목", "결과"],
        ["생산량 음수 검증", "0건 (통과)"],
        ["발주량 음수 검증", "0건 (통과)"],
        ["일정 정합성 (발주일 ≤ 필요일)", "100% 통과"],
        ["JSON 포맷 유효성", "100% 통과"],
    ]
    add_table_slide(slide, Inches(6.5), Inches(6.2), integrity_data,
                    [Inches(3.0), Inches(2.5)])

    # ============================================================
    # SLIDE 11: 검증 방법론
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
                 "07+  검증 방법론 — 신뢰성 확보 전략",
                 font_size=26, color=WHITE, bold=True)

    # 검증 방법론 3단계
    verify_cards = [
        ("1단계\n모델 검증", "Walk-Forward Cross Validation",
         "• 시계열 특성을 반영한 전진 교차검증\n• 주간 3-fold / 월간 2-fold\n• Pinball Loss 기반 분위 예측 정확도\n• Grid Search 하이퍼파라미터 튜닝\n• 예측 vs 실적 비교 리포트",
         MID_BLUE),
        ("2단계\n로직 검증", "리스크 스코어링 정합성",
         "• 등급 경계 연속성 검증 (갭 버그 수정)\n• 가중치 민감도 분석 (3개 시나리오)\n  → 모든 시나리오 C이상 11~16% 안정\n• 수정 전/후 비교: critical 66%→3%\n• 조치 큐 심각도 분포 정상화 확인",
         ACCENT),
        ("3단계\n결과 검증", "생산·발주 최적화 무결성",
         "• 음수 값 발생 여부 (0건 확인)\n• 일정 정합성 (발주일 ≤ 필요일)\n• JSON 포맷 유효성\n• 우선순위·긴급도 분포 합리성\n• 공급사 매핑율 확인\n• planned_qty=0 & critical 이상치 추적",
         GREEN),
    ]
    for i, (phase, subtitle, desc, color) in enumerate(verify_cards):
        x = Inches(0.5) + Inches(4.2) * i
        add_rect(slide, x, Inches(1.4), Inches(3.8), Inches(5.0), LIGHT_GRAY)
        add_rect(slide, x, Inches(1.4), Inches(3.8), Inches(1.2), color)
        add_text_box(slide, x + Inches(0.2), Inches(1.5), Inches(3.4), Inches(0.7),
                     phase, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(2.2), Inches(3.4), Inches(0.3),
                     subtitle, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.3), Inches(2.8), Inches(3.2), Inches(3.2),
                     desc, font_size=11, color=DARK_GRAY)

    # 하단: 지속적 검증 전략
    add_rect(slide, Inches(0.5), Inches(6.6), Inches(12), Inches(0.6), NAVY)
    add_text_box(slide, Inches(0.8), Inches(6.65), Inches(11.5), Inches(0.5),
                 "지속적 검증: 예측 vs 실적 자동 비교 리포트  |  주간 모델 재학습  |  리스크 스코어 추이 모니터링  |  A/B 테스트",
                 font_size=13, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 12: ERP 대비 차별점
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
                 "08  기대 효과 — ERP 대비 차별점 & 정량 목표",
                 font_size=26, color=WHITE, bold=True)

    # 비교 테이블
    comp_data = [
        ["구분", "기존 ERP", "본 솔루션"],
        ["과거 실적 조회", "O", "O"],
        ["미래 수요 예측", "X", "O (P10/P50/P90 확률 밴드)"],
        ["외부 변수 반영", "X", "O (30+ 경제·산업 지표)"],
        ["리스크 자동 감지", "X", "O (4유형, A~F 등급)"],
        ["생산 대응안 제안", "X", "O (자동 조치 큐)"],
        ["최적 생산량 산출", "X", "O (캐파+리스크 기반)"],
        ["발주 자동 추천", "X", "O (BOM+EOQ+공급사 추천)"],
        ["다차원 집계 분석", "제한적", "O (주별/월별/거래처별)"],
    ]
    add_table_slide(slide, Inches(0.5), Inches(1.3), comp_data,
                    [Inches(2.5), Inches(2.0), Inches(5.0)])

    # 정량 목표
    add_text_box(slide, Inches(0.5), Inches(5.3), Inches(5), Inches(0.4),
                 "정량적 기대 효과", font_size=18, color=NAVY, bold=True)
    effects = [
        ("20%↓", "재고 과잉 감소\n과잉 생산 방지", RED),
        ("80%↓", "엑셀 분석 시간\n절감", ACCENT),
        ("50%↓", "생산회의 준비\n시간 단축", MID_BLUE),
        ("사전 방지", "결품 납기 지연\n리스크 감지", GREEN),
    ]
    for i, (value, desc, color) in enumerate(effects):
        x = Inches(0.5) + Inches(3.15) * i
        add_rect(slide, x, Inches(5.8), Inches(2.9), Inches(1.4), color)
        add_text_box(slide, x + Inches(0.2), Inches(5.9), Inches(2.5), Inches(0.6),
                     value, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(6.5), Inches(2.5), Inches(0.5),
                     desc, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # SLIDE 13: 구축 로드맵
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
                 "09  구축 로드맵 — 추진 일정",
                 font_size=26, color=WHITE, bold=True)

    phases = [
        ("Phase 1", "데이터 탐색·전처리\nDB 구축", "완료",
         "• ERP CSV 9개 파일 적재\n• 27테이블 DB 구축\n• 외부 API 5종 연동\n• 환율·산업지수 17,000건",
         GREEN),
        ("Phase 2", "수요 변동성 분석\n모델 개발", "완료",
         "• 주간+월간 피처 엔지니어링\n  (46+35 피처)\n• LightGBM Quantile 학습\n• Walk-Forward CV\n• 이동평균 Fallback",
         GREEN),
        ("Phase 3", "재고 리스크\n점수화 엔진", "완료",
         "• 4유형 리스크 스코어링\n• A~F 등급 체계\n• 자동 조치 큐 생성\n• 가중치 민감도 분석\n• 실데이터 검증 완료",
         GREEN),
        ("Phase 4", "생산·발주\n최적화 알고리즘", "완료",
         "• 최적 생산량 산출 (S7)\n• BOM+EOQ 발주추천 (S8)\n• 공급사 가중점수 추천\n• 무결성 검증 통과",
         GREEN),
        ("Phase 5", "웹 대시보드\nAPI 서버", "진행 예정",
         "• FastAPI 백엔드 구축\n• Next.js 대시보드\n• Supabase Auth + RBAC\n• 실시간 알림 시스템",
         ORANGE),
        ("Phase 6", "통합 테스트\n배포", "진행 예정",
         "• E2E 테스트\n• 성능 최적화\n• Vercel + Supabase Cloud\n• 사용자 교육·매뉴얼",
         ORANGE),
    ]
    for i, (name, subtitle, status, desc, color) in enumerate(phases):
        x = Inches(0.3) + Inches(2.15) * i
        # 카드 배경
        add_rect(slide, x, Inches(1.4), Inches(2.0), Inches(5.5), LIGHT_GRAY)
        # 헤더
        add_rect(slide, x, Inches(1.4), Inches(2.0), Inches(1.3), color)
        add_text_box(slide, x + Inches(0.15), Inches(1.45), Inches(1.7), Inches(0.35),
                     name, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), Inches(1.85), Inches(1.7), Inches(0.5),
                     subtitle, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
        # 상태 배지
        badge_color = GREEN if status == "완료" else ORANGE
        add_rect(slide, x + Inches(0.35), Inches(2.45), Inches(1.3), Inches(0.28), badge_color)
        add_text_box(slide, x + Inches(0.35), Inches(2.45), Inches(1.3), Inches(0.28),
                     status, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # 설명
        add_text_box(slide, x + Inches(0.15), Inches(2.9), Inches(1.7), Inches(3.8),
                     desc, font_size=9, color=DARK_GRAY)

    # ============================================================
    # SLIDE 14: 리스크 및 대응 전략
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
                 "09+  리스크 및 대응 전략",
                 font_size=26, color=WHITE, bold=True)

    risk_table = [
        ["리스크", "영향", "대응 전략"],
        ["고객사별 패턴 편차", "예측 정확도 저하",
         "제품군 우선 적용 후 점진 확장, 고객사별 모델 분리"],
        ["외부 변수 영향도 불확실", "과적합/과소적합",
         "변수 자동 가중치 학습 (Feature Importance), 정기 재학습"],
        ["현업 미반영 가능성", "ROI 미달",
         'KPI 연계 리포트, "예측 vs 실제" 비교 대시보드'],
        ["ECOS API 연동 불안정", "한국 경제지표 결손",
         "ECOS 실데이터 + BOK 시뮬레이션 이중 적재"],
        ["LightGBM 미설치 환경", "모델 실행 불가",
         "이동평균 fallback 자동 적용"],
        ["공급사 매핑율 낮음 (11.7%)", "발주 추천 정확도 제한",
         "공급사 마스터 확충, 수동 매핑 병행"],
    ]
    add_table_slide(slide, Inches(0.5), Inches(1.4), risk_table,
                    [Inches(2.5), Inches(2.5), Inches(7.0)])

    # RBAC 권한 체계
    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(5), Inches(0.4),
                 "RBAC 권한 체계 (4단계)", font_size=16, color=NAVY, bold=True)
    rbac_data = [
        ["역할", "대시보드", "데이터 분석", "모델 설정", "보고서", "사용자 관리"],
        ["admin", "O", "O", "O", "O", "O"],
        ["manager", "O", "O", "O", "O", "-"],
        ["analyst", "O", "O", "-", "-", "-"],
        ["viewer", "O", "-", "-", "-", "-"],
    ]
    add_table_slide(slide, Inches(0.5), Inches(5.4), rbac_data,
                    [Inches(1.5), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3)])

    # ============================================================
    # SLIDE 15: Q&A / 마무리
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08), ACCENT)
    # 좌측 액센트 라인
    add_rect(slide, Inches(0.8), Inches(2.5), Inches(0.08), Inches(2.5), ACCENT)

    add_text_box(slide, Inches(1.2), Inches(2.5), Inches(8), Inches(0.8),
                 "Q & A",
                 font_size=48, color=WHITE, bold=True)
    add_rect(slide, Inches(1.2), Inches(3.5), Inches(3), Inches(0.04), ACCENT)
    add_text_box(slide, Inches(1.2), Inches(3.8), Inches(10), Inches(0.5),
                 "AI 기반 수요 변동성 최적화 SaaS",
                 font_size=22, color=ACCENT)
    add_text_box(slide, Inches(1.2), Inches(4.5), Inches(10), Inches(0.5),
                 "기존 ERP 데이터만으로 시작 가능  |  3개월 내 PoC  |  Phase 1~4 완료",
                 font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC))

    # 저장
    output_path = "c:/AI/PJTFNL/AI_수요예측_SaaS_제안서.pptx"
    prs.save(output_path)
    print(f"제안서 생성 완료: {output_path}")
    return output_path


if __name__ == "__main__":
    build_ppt()
